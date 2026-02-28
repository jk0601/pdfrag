"""
pptx_processor.py - PowerPoint 파일 처리기
============================================
PPTX 파일에서 텍스트와 이미지를 추출합니다.

[초보자 안내]
- python-pptx: PowerPoint(.pptx) 파일을 읽고 쓰는 Python 라이브러리
- PPTX 파일 구조:
  프레젠테이션 → 슬라이드들 → 각 슬라이드의 도형들(텍스트박스, 이미지, 표 등)
- 이 프로세서는:
  1. 각 슬라이드를 순회하며
  2. 텍스트박스, 표, 이미지 등 모든 도형에서 내용을 추출합니다

[처리 흐름]
  PPTX 파일 → 슬라이드별 순회 → 도형별 텍스트/이미지 추출 → 결합
"""

import os
import io
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

from processors.image_processor import ImageProcessor


@dataclass
class SlideContent:
    """한 슬라이드에서 추출된 내용"""

    slide_number: int
    title: str = ""
    texts: list[str] = field(default_factory=list)
    table_texts: list[str] = field(default_factory=list)
    image_texts: list[str] = field(default_factory=list)

    @property
    def full_text(self) -> str:
        """슬라이드의 전체 내용을 하나의 텍스트로 결합"""
        parts = []
        if self.title:
            parts.append(f"[제목] {self.title}")
        parts.extend(self.texts)
        for table_text in self.table_texts:
            parts.append(f"[표]\n{table_text}")
        for img_text in self.image_texts:
            if img_text.strip():
                parts.append(f"[이미지 내용]: {img_text}")
        return "\n".join(parts)


class PPTXProcessor:
    """
    PPTX 파일에서 텍스트와 이미지를 추출하는 프로세서

    사용 예시:
        processor = PPTXProcessor()
        result = processor.process("발표자료.pptx")
        for slide in result["slides"]:
            print(f"슬라이드 {slide.slide_number}: {slide.full_text}")
    """

    def __init__(self, ocr_enabled: bool = True):
        self.ocr_enabled = ocr_enabled
        self.image_processor = ImageProcessor() if ocr_enabled else None

    def process(self, file_path: str) -> dict:
        """
        PPTX 파일을 처리하여 텍스트를 추출합니다.

        Returns:
            dict: {
                "filename": 파일 이름,
                "file_type": "pptx",
                "file_size": 파일 크기,
                "page_count": 슬라이드 수,
                "slides": [SlideContent, ...],
                "full_text": 전체 텍스트,
            }
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"파일을 찾을 수 없습니다: {file_path}")

        prs = Presentation(file_path)
        slides: list[SlideContent] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_content = self._extract_slide(slide, slide_num)
            slides.append(slide_content)

        full_text = "\n\n".join(
            f"--- 슬라이드 {s.slide_number} ---\n{s.full_text}" for s in slides
        )

        return {
            "filename": os.path.basename(file_path),
            "file_type": "pptx",
            "file_size": os.path.getsize(file_path),
            "page_count": len(slides),
            "slides": slides,
            "full_text": full_text,
        }

    def _extract_slide(self, slide, slide_number: int) -> SlideContent:
        """한 슬라이드에서 모든 내용을 추출합니다."""
        content = SlideContent(slide_number=slide_number)

        # 슬라이드 제목 추출
        if slide.shapes.title:
            content.title = slide.shapes.title.text.strip()

        for shape in slide.shapes:
            # 텍스트 프레임이 있는 도형 (텍스트박스 등)
            if shape.has_text_frame:
                text = self._extract_text_frame(shape.text_frame)
                if text and text != content.title:
                    content.texts.append(text)

            # 표(Table)
            if shape.has_table:
                table_text = self._extract_table(shape.table)
                if table_text:
                    content.table_texts.append(table_text)

            # 이미지
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if self.ocr_enabled and self.image_processor:
                    img_text = self._extract_image(shape)
                    if img_text:
                        content.image_texts.append(img_text)

            # 그룹 도형 (여러 도형이 묶인 것)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self._extract_group(shape, content)

        return content

    def _extract_text_frame(self, text_frame) -> str:
        """텍스트 프레임에서 텍스트를 추출합니다."""
        paragraphs = []
        for paragraph in text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                paragraphs.append(text)
        return "\n".join(paragraphs)

    def _extract_table(self, table) -> str:
        """표에서 텍스트를 추출합니다."""
        rows = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cells.append(cell.text.strip())
            rows.append(" | ".join(cells))
        return "\n".join(rows)

    def _extract_image(self, shape) -> str:
        """이미지 도형에서 OCR로 텍스트를 추출합니다."""
        try:
            image_blob = shape.image.blob
            img = Image.open(io.BytesIO(image_blob))
            return self.image_processor.extract_text_from_image(img)
        except Exception:
            return ""

    def _extract_group(self, group_shape, content: SlideContent):
        """그룹 도형 내부의 텍스트를 재귀적으로 추출합니다."""
        for shape in group_shape.shapes:
            if shape.has_text_frame:
                text = self._extract_text_frame(shape.text_frame)
                if text:
                    content.texts.append(text)
            if shape.has_table:
                table_text = self._extract_table(shape.table)
                if table_text:
                    content.table_texts.append(table_text)
